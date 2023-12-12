import os
import pathlib
from datetime import date, datetime

import pptx
import pptx.presentation

from typing import Any, Dict, Optional, Union

try:
    from plugins.com import powerpoint
    IMPORT_PYWIN32COM_OK = True
except ModuleNotFoundError:
    IMPORT_PYWIN32COM_OK = False

class PowerPointRecipes:

    @classmethod
    def read_all(cls, presentation: Union[str, pptx.presentation.Presentation]) -> str:
        """
        Extrai o conteúdo de texto de todo o documento.

        :param document: pode ser o caminho do documento ou o objeto do documento
        :returns: string contendo todo o texto do documento
        """
        assert isinstance(presentation, (str, pptx.presentation.Presentation)), 'Tipo inválido de documento'

        if isinstance(presentation, str):
            ppt = pptx.Presentation(presentation)
        else:
            ppt = presentation

        # percorre todos os elementos da apresentação
        paragraphs = []
        for slide in ppt.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    paragraphs.extend(shape.text_frame.paragraphs)

        # agora junta tudo e retorna
        return '\n'.join([p.text for p in paragraphs])

    @classmethod
    def fill_template(cls, 
                      filename: str, 
                      save_as: str,
                      values_dict: Dict[str, Any], 
                      *, 
                      placeholder_pattern: Optional[str] = '{{placeholder}}'):
        """
        Preenche uma apresentação substituindo termos por valores de um dicionário.

        :param filename: nome do arquivo de entrada
        :param save_as: nome do arquivo de saída
        :param values_dict: dicionário contendo os termos e os valores a serem substituídos
        :param placeholder_pattern: padrão do placeholder. Note que a palavra "placeholder" deve constar no padrão.
        """
        assert filename.lower().endswith('.pptx'), 'Arquivo de entrada precisa ser do tipo PPTX'
        assert save_as.lower().endswith('.pptx'), 'Arquivo de saída precisa ser do tipo PPTX'
        assert filename.lower()!=save_as.lower(), 'Os arquivos de entrada e saída não podem ser os mesmos.'
        assert 'placeholder' in placeholder_pattern, 'O padrão de substituição precisa conter a palavra "placeholder"'

        known_converters = {
            float: lambda v: str(v).replace('.', ','),
            date: lambda v: v.strftime(r'%d/%m/%Y'),
            datetime: lambda v: v.strftime(r'%d/%m/%Y'),
        }

        presentation = pptx.Presentation(filename)

        paragraphs = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    paragraphs.extend(shape.text_frame.paragraphs)

        for p in paragraphs:
            for k, v in values_dict.items():
                placeholder = placeholder_pattern.replace('placeholder', str(k))
                
                v = v or ''
                converted_v = known_converters.get(type(v), str)(v)

                if placeholder in p.text:
                    inline = p.runs
                    
                    # esse tipo de documento é complexo, podemos ter um placeholder distribuído em mais de um run
                    started = False
                    key_index = 0
                    # found_runs é uma lista com (run index, index of match, length of match)
                    found_runs = []
                    found_all = False
                    replace_done = False

                    for i, run in enumerate(inline):

                        # caso 1: encontrou no run inteiro
                        if placeholder in run.text and not started:
                            found_runs.append((i, run.text.find(placeholder), len(placeholder)))
                            text = run.text.replace(placeholder, converted_v)
                            run.text = text
                            replace_done = True
                            found_all = True
                            break

                        if placeholder[key_index] not in run.text and not started:
                            # continua procurando...
                            continue

                        # caso 2: procurar por texto parcial, primeira parte
                        if placeholder[key_index] in run.text and run.text[-1] in placeholder and not started:
                            # verifica sequencia
                            start_index = run.text.find(placeholder[key_index])
                            check_length = len(run.text)
                            for text_index in range(start_index, check_length):
                                if run.text[text_index] != placeholder[key_index]:
                                    # falso positivo
                                    break
                            if key_index == 0:
                                started = True
                            chars_found = check_length - start_index
                            key_index += chars_found
                            found_runs.append((i, start_index, chars_found))
                            if key_index != len(placeholder):
                                continue
                            else:
                                # todos os caracteres foram encontrados
                                found_all = True
                                break

                        # caso 2: procurar por texto parcial, segunda parte
                        if placeholder[key_index] in run.text and started and not found_all:
                            # verifica sequencia
                            chars_found = 0
                            check_length = len(run.text)
                            for text_index in range(0, check_length):
                                if run.text[text_index] == placeholder[key_index]:
                                    key_index += 1
                                    chars_found += 1
                                else:
                                    break
                            # encontrou
                            found_runs.append((i, 0, chars_found))
                            if key_index == len(placeholder):
                                found_all = True
                                break

                    # replace dos runs distribuidos
                    if found_all and not replace_done:
                        for i, item in enumerate(found_runs):
                            index, start, length = [t for t in item]
                            if i == 0:
                                text = inline[index].text.replace(inline[index].text[start:start + length], converted_v)
                                inline[index].text = text
                            else:
                                text = inline[index].text.replace(inline[index].text[start:start + length], '')
                                inline[index].text = text

        if os.path.exists(save_as):
            os.remove(save_as)

        presentation.save(save_as)

    @classmethod
    def to_pdf(cls, pp_filename: str, pdf_filename: str):
        """
        Converte arquivo PPTX em PDF. Necessário ter o PowerPoint instalado.

        :param pp_filename: nome do arquivo PPTX de entrada
        :param pdf_filename: nome do arquivo PDF de saída
        """
        assert IMPORT_PYWIN32COM_OK, 'Os requisitos para essa funcionalidade não foram instalados/ configurados.'
        
        with powerpoint.application(kill_after=True) as pp:
            pp_filename = str(pathlib.Path(pp_filename).resolve())
            pdf_filename = str(pathlib.Path(pdf_filename).resolve())

            ppt = pp.Presentations.Open(pp_filename, WithWindow=0)

            if os.path.exists(pdf_filename):
                os.remove(pdf_filename)

            ppt.SaveAs(pdf_filename, FileFormat=powerpoint.constants.ppSaveAsPDF)
            ppt.Close()
        return